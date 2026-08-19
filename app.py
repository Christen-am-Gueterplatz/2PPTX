# app.py
import io
import os
from flask import Flask, request, render_template, send_file, jsonify, session, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps # ImageOps für EXIF-Orientierung hinzufügen
import fitz  # PyMuPDF
import traceback # Für detaillierteres Error-Logging
import gc  # Garbage Collector
from dotenv import load_dotenv
from werkzeug.exceptions import RequestEntityTooLarge

# Lade Umgebungsvariablen aus .env Datei
load_dotenv()

# --- Konstanten ---
SLIDE_WIDTH_EMU = 9144000 # 10 inches
SLIDE_HEIGHT_EMU = 5143500 # 5.625 inches
MAX_IMG_WIDTH_PX = 1920
MAX_IMG_HEIGHT_PX = 1080

# --- Flask App Initialisierung ---
app = Flask(__name__)
# Erhöhe das Limit für potenziell viele Dateien oder große PDFs (Standard: 500 MB)
MAX_UPLOAD_MB = int(os.getenv('MAX_UPLOAD_MB', '500'))
app.config['MAX_CONTENT_LENGTH'] = MAX_UPLOAD_MB * 1024 * 1024
# Secret Key für Sessions (sollte in Produktion ein zufälliger String sein)
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'dev-secret-key-change-in-production')
# Passwort aus Umgebungsvariable
APP_PASSWORD = os.getenv('APP_PASSWORD', 'password')

# --- NEUE Hilfsfunktion: Bild verarbeiten (Größe, Orientierung) ---
def process_image_stream(image_stream):
    """
    Öffnet ein Bild aus einem Stream, korrigiert die Orientierung basierend auf EXIF,
    ändert die Größe, falls es 1920x1080 überschreitet (unter Beibehaltung des Seitenverhältnisses),
    und gibt einen neuen Stream mit dem verarbeiteten Bild (als PNG) zurück.
    """
    try:
        img = Image.open(image_stream)

        # 1. EXIF-Orientierung korrigieren
        #    ImageOps.exif_transpose liest EXIF und wendet Rotation/Spiegelung an
        img = ImageOps.exif_transpose(img)
        print(f"EXIF-Orientierung angewendet (falls vorhanden).")

        # 2. Größe ändern, wenn nötig
        current_width, current_height = img.size
        print(f"Originalgröße: {current_width}x{current_height}")

        if current_width > MAX_IMG_WIDTH_PX or current_height > MAX_IMG_HEIGHT_PX:
            # Seitenverhältnis beibehalten, auf Maximalgröße reduzieren
            # thumbnail ändert das Bild Objekt 'in-place'
            img.thumbnail((MAX_IMG_WIDTH_PX, MAX_IMG_HEIGHT_PX), Image.Resampling.LANCZOS)
            new_width, new_height = img.size
            print(f"Größe geändert auf: {new_width}x{new_height}")
        else:
            print("Keine Größenänderung erforderlich.")

        # 3. Verarbeitetes Bild in einen neuen Stream speichern (PNG für verlustfreie Zwischenspeicherung)
        output_stream = io.BytesIO()
        img.save(output_stream, format='PNG')
        output_stream.seek(0) # Wichtig: Stream zurücksetzen
        img.close() # Original PIL Image schließen
        return output_stream

    except Exception as e:
        print(f"Fehler beim Verarbeiten des Bild-Streams:")
        traceback.print_exc()
        # Im Fehlerfall versuchen, den Originalstream zurückzugeben (oder None)
        image_stream.seek(0) # Zurücksetzen für den Fall, dass es doch noch geht
        return image_stream # Oder return None und später darauf prüfen


# --- Hilfsfunktion: Bild zentriert auf Folie hinzufügen ---
def add_image_centered(slide, image_stream):
    """Fügt ein Bild aus einem Stream zentriert auf einer Folie hinzu."""
    try:
        image_stream.seek(0)
        with Image.open(image_stream) as img:
            img_width_px, img_height_px = img.size

        image_stream.seek(0)

        # Skalierungsfaktor für Folie berechnen
        scale_w = float(SLIDE_WIDTH_EMU) / img_width_px
        scale_h = float(SLIDE_HEIGHT_EMU) / img_height_px
        scale = min(scale_w, scale_h)

        pic_width_emu = int(img_width_px * scale)
        pic_height_emu = int(img_height_px * scale)
        left = int((SLIDE_WIDTH_EMU - pic_width_emu) / 2)
        top = int((SLIDE_HEIGHT_EMU - pic_height_emu) / 2)

        slide.shapes.add_picture(image_stream, left, top, width=pic_width_emu, height=pic_height_emu)
    except Exception as e:
        print(f"Fehler beim Hinzufügen des Bildes zur Folie: {e}")
        traceback.print_exc()
        try:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tf = textbox.text_frame
            tf.text = f"Fehler beim Laden/Platzieren eines Bildes:\n{e}"
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            tf.word_wrap = True
        except Exception:
            pass

def add_slide_with_image(prs, image_stream):
    """Erstellt eine schwarze Folie und platziert das Bild zentriert."""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    add_image_centered(slide, image_stream)

# --- Route: /upload (geschützt) ---
@app.route('/upload', methods=['POST'])
def upload_files():
    if not check_auth():
        return jsonify({"error": "Nicht authentifiziert"}), 401
    
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
    slides_count = 0

    try:
        files = request.files.getlist('files')
        if not files or all(f.filename == '' for f in files):
            return jsonify({"error": "Keine Dateien ausgewählt"}), 400

        for file in files:
            try:
                filename = file.filename.lower()
                print(f"Verarbeite Datei: {file.filename}")
                file_bytes = file.read()

                if filename.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')):
                    stream = io.BytesIO(file_bytes)
                    processed_stream = process_image_stream(stream)
                    stream.close()
                    if processed_stream:
                        add_slide_with_image(prs, processed_stream)
                        processed_stream.close()
                        slides_count += 1
                
                elif filename.endswith('.pdf'):
                    pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
                    try:
                        for page_num in range(len(pdf_document)):
                            page = pdf_document.load_page(page_num)
                            rect = page.rect
                            
                            # Native C-Skalierung direkt auf max 1920x1080 berechnen
                            if rect.width > 0 and rect.height > 0:
                                scale = min(MAX_IMG_WIDTH_PX / rect.width, MAX_IMG_HEIGHT_PX / rect.height)
                                matrix = fitz.Matrix(scale, scale)
                            else:
                                matrix = fitz.Matrix(2.0, 2.0)

                            # Direkt als JPEG rendern (extrem schnell, geringer RAM-Bedarf, hohe Qualität)
                            pix = page.get_pixmap(matrix=matrix, alpha=False)
                            img_bytes = pix.tobytes("jpeg", jpg_quality=92)
                            page_stream = io.BytesIO(img_bytes)
                            
                            add_slide_with_image(prs, page_stream)
                            page_stream.close()
                            del pix
                            slides_count += 1
                    finally:
                        pdf_document.close()
                
                del file_bytes
                gc.collect()
            
            except Exception as e:
                print(f"Fehler bei Datei {file.filename}: {str(e)}")
                traceback.print_exc()
                continue

        if slides_count == 0:
            return jsonify({"error": "Keine gültigen Seiten/Bilder gefunden"}), 400

        # Präsentation in Byte-Stream speichern
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        del prs
        gc.collect()

        return send_file(
            pptx_io,
            as_attachment=True,
            download_name='presentation.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except RequestEntityTooLarge:
        del prs
        gc.collect()
        return jsonify({"error": f"Datei(en) zu groß! Das maximale Upload-Limit liegt bei {MAX_UPLOAD_MB} MB."}), 413

    except Exception as e:
        del prs
        gc.collect()
        print(f"Verarbeitungsfehler: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Verarbeitungsfehler: {str(e)}"}), 500

@app.errorhandler(413)
def request_entity_too_large(error):
    return jsonify({"error": f"Datei(en) zu groß! Das maximale Upload-Limit liegt bei {MAX_UPLOAD_MB} MB."}), 413

# --- Authentifizierungs-Middleware ---
def check_auth():
    """Prüft, ob der Benutzer authentifiziert ist."""
    return session.get('authenticated', False)

# --- Login Route ---
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        password = request.form.get('password', '')
        if password == APP_PASSWORD:
            session['authenticated'] = True
            return redirect(url_for('index'))
        else:
            return render_template('login.html', error='Falsches Passwort')
    return render_template('login.html')

# --- Logout Route ---
@app.route('/logout')
def logout():
    session.pop('authenticated', None)
    return redirect(url_for('login'))

# --- Index Route (geschützt) ---
@app.route('/')
def index():
    if not check_auth():
        return redirect(url_for('login'))
    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=False, host='0.0.0.0', port=5000)

